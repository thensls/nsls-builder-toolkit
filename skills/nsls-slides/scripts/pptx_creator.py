#!/usr/bin/env python3
"""
NSLS/Society brand PowerPoint creator.

Usage:
  echo '{"slides":[...]}' | python3 pptx_creator.py --output /tmp/out.pptx
  python3 pptx_creator.py --input slides.json --output /tmp/out.pptx

Input JSON schema:
  {
    "slides": [
      {
        "layout": "title",
        "headline": "Find your people.\nFind your path.",
        "subhead": "The community-driven success platform."
      },
      {
        "layout": "section",
        "text": "Our Mission",
        "bg": "yellow"
      },
      {
        "layout": "content",
        "title": "What We Do",
        "body": "Optional paragraph of body text.",
        "bullets": ["Point one", "Point two", "Point three"]
      },
      {
        "layout": "two_column",
        "title": "Comparison",
        "left": "Left column text",
        "right": "Right column text"
      },
      {
        "layout": "quote",
        "text": "Leadership is not about a title.",
        "attribution": "— John Maxwell",
        "bg": "lavender"
      }
    ]
  }

Background color options: cream, yellow, lavender, pink, green, taupe, espresso
"""

import io
import os
import sys
import json
import zipfile
import argparse
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Brand Tokens ────────────────────────────────────────────────────────────

# --- Society brand (by the NSLS) — Cigars + Inter, cream/yellow palette ---
SOCIETY_COLORS = {
    "cream":    RGBColor(0xFA, 0xF8, 0xEE),  # #FAF8EE — primary background
    "espresso": RGBColor(0x20, 0x14, 0x14),  # #201414 — primary text / dark bg
    "yellow":   RGBColor(0xF2, 0xDA, 0x4E),  # #F2DA4E
    "lavender": RGBColor(0x96, 0x9B, 0xDE),  # #969BDE
    "pink":     RGBColor(0xF3, 0xAE, 0xE6),  # #F3AEE6
    "green":    RGBColor(0x9B, 0xD7, 0x78),  # #9BD778
    "taupe":    RGBColor(0xC8, 0xBD, 0xAF),  # #C8BDAF
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
}

SOCIETY_DARK_BGS = {"espresso"}

SOCIETY_FONTS = {
    "headline": "HW Cigars Medium",
    "logotype": "HW Cigars SemiBold",
    "body":     "Inter",
}

SOCIETY_EMBED_FONTS = [
    {
        "typeface": "HW Cigars Medium",
        "style":    "regular",
        "path":     os.path.expanduser("~/Library/Fonts/HW Cigars Medium.otf"),
    },
    {
        "typeface": "HW Cigars SemiBold",
        "style":    "regular",
        "path":     os.path.expanduser("~/Library/Fonts/HW Cigars SemiBold.otf"),
    },
]

# --- NSLS brand (honor society) — Lexend Deca + Avenir, navy/teal/gold palette ---
NSLS_COLORS = {
    "white":     RGBColor(0xFF, 0xFF, 0xFF),  # #FFFFFF — primary background
    "navy":      RGBColor(0x18, 0x31, 0x5A),  # #18315A — dark bg / accent
    "darkblue":  RGBColor(0x42, 0x5B, 0x76),  # #425B76 — secondary text / bg
    "bluegray":  RGBColor(0x33, 0x47, 0x5B),  # #33475B — primary text
    "teal":      RGBColor(0x00, 0x91, 0xAE),  # #0091AE — links / CTA accent
    "gold":      RGBColor(0xEE, 0xB1, 0x17),  # #EEB117 — honor society gold
    "lightblue": RGBColor(0xE5, 0xF5, 0xF8),  # #E5F5F8 — light accent bg
    "nearblack": RGBColor(0x19, 0x19, 0x19),  # #191919 — near-black text
}

NSLS_DARK_BGS = {"navy", "darkblue", "bluegray", "nearblack"}

NSLS_FONTS = {
    "headline": "Lexend Deca",
    "logotype": "Lexend Deca",
    "body":     "Avenir",
}

NSLS_EMBED_FONTS = []  # Lexend Deca and Avenir are system/Google fonts — no embedding needed

# --- Active brand (set by --brand flag, defaults to Society) ---
COLORS = SOCIETY_COLORS
DARK_BGS = SOCIETY_DARK_BGS
FONTS = SOCIETY_FONTS
EMBED_FONTS = SOCIETY_EMBED_FONTS


def set_brand(brand_name):
    """Switch the active brand tokens. Called from main() based on --brand flag."""
    global COLORS, DARK_BGS, FONTS, EMBED_FONTS
    if brand_name == "nsls":
        COLORS = NSLS_COLORS
        DARK_BGS = NSLS_DARK_BGS
        FONTS = NSLS_FONTS
        EMBED_FONTS = NSLS_EMBED_FONTS
    elif brand_name == "society":
        COLORS = SOCIETY_COLORS
        DARK_BGS = SOCIETY_DARK_BGS
        FONTS = SOCIETY_FONTS
        EMBED_FONTS = SOCIETY_EMBED_FONTS
    else:
        raise ValueError(f"Unknown brand: {brand_name!r}. Use 'nsls' or 'society'.")

# OOXML namespace constants for font embedding
_PRES_NS     = "http://schemas.openxmlformats.org/presentationml/2006/main"
_REL_NS      = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_FONT_REL    = f"{_REL_NS}/font"
_RELS_FILE   = "ppt/_rels/presentation.xml.rels"
_PRES_FILE   = "ppt/presentation.xml"

# Slide dimensions: 16:9 widescreen (matches Google Slides default)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Layout margins
MH = Inches(0.9)               # horizontal margin
MV = Inches(0.7)               # vertical margin (top)
CW = SLIDE_W - 2 * MH         # content width
RULE_H = Pt(1.5)               # thin yellow rule height


# ── Font Embedding ───────────────────────────────────────────────────────────

def embed_fonts(pptx_path, font_specs=None):
    """
    Post-process a saved PPTX to embed font binaries per OOXML spec.

    Modifies the file in-place. Silently skips fonts whose file path doesn't
    exist (e.g. running on a machine without HW Cigars installed).
    """
    if font_specs is None:
        font_specs = EMBED_FONTS

    specs = [s for s in font_specs if os.path.exists(s["path"])]
    if not specs:
        return  # nothing to embed

    # Assign stable internal IDs and zip paths
    entries = [
        {**s, "rId": f"rIdFont{i + 1}", "zip_path": f"ppt/fonts/font{i + 1}.fntdata"}
        for i, s in enumerate(specs)
    ]

    src = Path(pptx_path).read_bytes()
    out = io.BytesIO()

    with zipfile.ZipFile(io.BytesIO(src), "r") as zin, \
         zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == _RELS_FILE:
                data = _patch_rels(data, entries)
            elif item.filename == _PRES_FILE:
                data = _patch_pres_xml(data, entries)

            zout.writestr(item, data)

        # Append font binary files
        for entry in entries:
            zout.writestr(entry["zip_path"], Path(entry["path"]).read_bytes())

    Path(pptx_path).write_bytes(out.getvalue())


def _patch_rels(data, entries):
    """Insert font <Relationship> entries into presentation rels XML."""
    new_rels = "\n".join(
        f'  <Relationship Id="{e["rId"]}" Type="{_FONT_REL}" Target="fonts/font{i + 1}.fntdata"/>'
        for i, e in enumerate(entries)
    )
    return data.decode("utf-8").replace(
        "</Relationships>", f"{new_rels}\n</Relationships>"
    ).encode("utf-8")


def _patch_pres_xml(data, entries):
    """Add <p:embeddedFontLst> block to presentation.xml."""
    root = etree.fromstring(data)

    ns  = _PRES_NS
    rns = _REL_NS

    # Find or create the font list element
    lst_tag  = f"{{{ns}}}embeddedFontLst"
    font_list = root.find(lst_tag)
    if font_list is None:
        font_list = etree.SubElement(root, lst_tag)

    # Group entries by typeface so each typeface gets one <p:embeddedFont>
    seen = {}
    for entry in entries:
        seen.setdefault(entry["typeface"], []).append(entry)

    for typeface, specs in seen.items():
        ef = etree.SubElement(font_list, f"{{{ns}}}embeddedFont")
        fi = etree.SubElement(ef, f"{{{ns}}}font")
        fi.set("typeface", typeface)
        for spec in specs:
            style_el = etree.SubElement(ef, f"{{{ns}}}{spec['style']}")
            style_el.set(f"{{{rns}}}id", spec["rId"])

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ── Core Helpers ────────────────────────────────────────────────────────────

def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _default_bg():
    """Return the brand's default background color name."""
    return "cream" if "cream" in COLORS else "white"


def _light_text():
    """Return the brand's light text color (for dark backgrounds)."""
    return COLORS.get("cream", COLORS.get("white"))


def _dark_text():
    """Return the brand's dark text color (for light backgrounds)."""
    return COLORS.get("espresso", COLORS.get("bluegray", COLORS.get("nearblack")))


def blank_slide(prs, bg=None):
    """Add a completely blank slide with a solid brand background."""
    if bg is None:
        bg = _default_bg()
    blank_layout = prs.slide_layouts[6]  # index 6 = Blank in default template
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS.get(bg, COLORS[_default_bg()])
    return slide


def text_color_for_bg(bg):
    """Return light text for dark backgrounds, dark text for light."""
    return _light_text() if bg in DARK_BGS else _dark_text()


def add_textbox(slide, left, top, width, height, text,
                font_name, size, color=None, align=PP_ALIGN.LEFT):
    """Add a single-paragraph text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = size
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullets(slide, left, top, width, height, items,
                font_name, size, color=None):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"\u2022  {text}"  # bullet character
        run.font.name = font_name
        run.font.size = size
        if color:
            run.font.color.rgb = color
    return txBox


def _accent_color():
    """Return the brand's accent color name (yellow for Society, gold for NSLS)."""
    return "yellow" if "yellow" in COLORS else "gold"


def add_rule(slide, left, top, width, color_key=None):
    """Add a thin horizontal rule (1.5pt tall colored rectangle)."""
    if color_key is None:
        color_key = _accent_color()
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, RULE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color_key]
    shape.line.fill.background()  # no border line


# ── Slide Builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    """
    Large Cigars Medium headline + Inter subhead on cream.

    Fields: headline (str), subhead (str, optional)
    """
    slide = blank_slide(prs)
    headline = data.get("headline", "")
    subhead = data.get("subhead", "")

    # Headline — upper 60% of slide
    add_textbox(
        slide,
        left=MH, top=Inches(0.9),
        width=CW, height=Inches(3.0),
        text=headline,
        font_name=FONTS["headline"],
        size=Pt(56),
        color=_dark_text(),
        align=PP_ALIGN.LEFT,
    )

    # Subhead
    if subhead:
        add_textbox(
            slide,
            left=MH, top=Inches(4.1),
            width=CW, height=Inches(0.9),
            text=subhead,
            font_name=FONTS["body"],
            size=Pt(20),
            color=_dark_text(),
            align=PP_ALIGN.LEFT,
        )

    return slide


def build_section_slide(prs, data):
    """
    Bold headline on brand color background — section divider.

    Fields: text (str), bg (color key, default "yellow")
    """
    bg = data.get("bg", _accent_color())
    slide = blank_slide(prs, bg=bg)
    text = data.get("text", "")

    add_textbox(
        slide,
        left=MH, top=Inches(1.8),
        width=CW, height=Inches(2.5),
        text=text,
        font_name=FONTS["headline"],
        size=Pt(56),
        color=text_color_for_bg(bg),
        align=PP_ALIGN.LEFT,
    )
    return slide


def build_content_slide(prs, data):
    """
    Title + optional body paragraph + optional bullet list on cream.

    Fields: title (str), body (str, optional), bullets (list[str], optional)
    """
    slide = blank_slide(prs)
    title = data.get("title", "")
    body = data.get("body", "")
    bullets = data.get("bullets", [])

    # Title
    add_textbox(
        slide,
        left=MH, top=MV,
        width=CW, height=Inches(0.8),
        text=title,
        font_name=FONTS["headline"],
        size=Pt(32),
        color=_dark_text(),
    )

    # Yellow rule below title
    add_rule(slide, left=MH, top=Inches(1.6), width=CW)

    content_top = Inches(1.85)
    content_h = Inches(3.2)

    if body:
        add_textbox(
            slide,
            left=MH, top=content_top,
            width=CW, height=Inches(1.2),
            text=body,
            font_name=FONTS["body"],
            size=Pt(17),
            color=_dark_text(),
        )
        content_top += Inches(1.3)
        content_h -= Inches(1.3)

    if bullets:
        add_bullets(
            slide,
            left=MH, top=content_top,
            width=CW, height=content_h,
            items=bullets,
            font_name=FONTS["body"],
            size=Pt(17),
            color=_dark_text(),
        )

    return slide


def build_two_column_slide(prs, data):
    """
    Title + two equal text columns on cream background.

    Fields: title (str), left (str), right (str)
    """
    slide = blank_slide(prs)
    title = data.get("title", "")
    left_text = data.get("left", "")
    right_text = data.get("right", "")

    gap = Inches(0.5)
    col_w = (CW - gap) / 2

    # Title
    add_textbox(
        slide,
        left=MH, top=MV,
        width=CW, height=Inches(0.8),
        text=title,
        font_name=FONTS["headline"],
        size=Pt(32),
        color=_dark_text(),
    )

    add_rule(slide, left=MH, top=Inches(1.6), width=CW)

    col_top = Inches(1.85)
    col_h = Inches(3.2)

    add_textbox(
        slide,
        left=MH, top=col_top,
        width=col_w, height=col_h,
        text=left_text,
        font_name=FONTS["body"],
        size=Pt(16),
        color=_dark_text(),
    )

    add_textbox(
        slide,
        left=MH + col_w + gap, top=col_top,
        width=col_w, height=col_h,
        text=right_text,
        font_name=FONTS["body"],
        size=Pt(16),
        color=_dark_text(),
    )

    return slide


def build_quote_slide(prs, data):
    """
    Large pull quote centered on brand color background.

    Fields: text (str), attribution (str, optional), bg (color key, default "yellow")
    """
    bg = data.get("bg", _accent_color())
    slide = blank_slide(prs, bg=bg)
    text = data.get("text", "")
    attribution = data.get("attribution", "")

    tc = text_color_for_bg(bg)

    add_textbox(
        slide,
        left=MH, top=Inches(1.0),
        width=CW, height=Inches(3.0),
        text=f"\u201c{text}\u201d",
        font_name=FONTS["headline"],
        size=Pt(40),
        color=tc,
        align=PP_ALIGN.CENTER,
    )

    if attribution:
        add_textbox(
            slide,
            left=MH, top=Inches(4.3),
            width=CW, height=Inches(0.7),
            text=attribution,
            font_name=FONTS["body"],
            size=Pt(14),
            color=tc,
            align=PP_ALIGN.CENTER,
        )

    return slide


BUILDERS = {
    "title":      build_title_slide,
    "section":    build_section_slide,
    "content":    build_content_slide,
    "two_column": build_two_column_slide,
    "quote":      build_quote_slide,
}


# ── Entry Point ─────────────────────────────────────────────────────────────

def build_presentation(data, output_path):
    prs = new_presentation()
    slides_data = data.get("slides", [])

    if not slides_data:
        print("Warning: no slides in input data", file=sys.stderr)

    for i, slide_data in enumerate(slides_data):
        layout = slide_data.get("layout", "content")
        builder = BUILDERS.get(layout)
        if not builder:
            print(f"Warning: unknown layout '{layout}' on slide {i + 1}, skipping", file=sys.stderr)
            continue
        builder(prs, slide_data)

    prs.save(output_path)
    embed_fonts(output_path)
    print(f"Saved: {output_path}")


def export_pdf(pptx_path, pdf_path):
    """
    Export a PPTX to PDF via Keynote (macOS only).

    Keynote imports the PPTX and exports a PDF with fonts properly embedded —
    HW Cigars renders correctly in the resulting PDF everywhere, including
    Google Drive's PDF viewer. Requires Keynote to be installed.
    """
    import subprocess
    script = f"""
tell application "Keynote"
    open POSIX file "{pptx_path}"
    delay 5
    export front document to POSIX file "{pdf_path}" as PDF
    close front document saving no
end tell
"""
    result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"Keynote PDF export failed: {result.stderr.strip()}")
    print(f"PDF:   {pdf_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Generate NSLS/Society-branded PowerPoint slides"
    )
    parser.add_argument("--input",  "-i", help="JSON input file (default: stdin)")
    parser.add_argument("--output", "-o", required=True, help="Output .pptx path")
    parser.add_argument("--brand", "-b", choices=["society", "nsls"], default="society",
                        help="Brand: 'society' (Cigars+Inter, cream/yellow) or 'nsls' (Lexend Deca+Avenir, navy/teal/gold)")
    parser.add_argument("--pdf", action="store_true",
                        help="Also export a PDF via Keynote (fonts render correctly everywhere)")
    args = parser.parse_args()

    set_brand(args.brand)

    if args.input:
        with open(args.input) as f:
            data = json.load(f)
    else:
        if sys.stdin.isatty():
            print("Reading JSON from stdin... (paste JSON and press Ctrl+D)", file=sys.stderr)
        data = json.load(sys.stdin)

    build_presentation(data, args.output)

    if args.pdf:
        pdf_path = str(Path(args.output).with_suffix(".pdf"))
        export_pdf(str(Path(args.output).resolve()), pdf_path)


if __name__ == "__main__":
    main()
